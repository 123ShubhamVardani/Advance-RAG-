# 🤖 LangChain AI Chatbot - Production Ready

import os
import streamlit as st
from dotenv import load_dotenv
from datetime import datetime
import hashlib
import json
import ssl
import traceback
from pathlib import Path
from logger import logger
from typing import Optional, Union, Any

# Try to use Windows certificate store for Requests (corporate CA support)
try:
    pass  # type: ignore
except Exception:
    pass

# Configure SSL only if explicitly allowed (prevents silently disabling verification)
if os.getenv("ALLOW_INSECURE_SSL", "false").lower() == "true":
    ssl._create_default_https_context = ssl._create_unverified_context
    try:
        logger.warning("Global SSL verification disabled via ALLOW_INSECURE_SSL")
    except Exception:
        pass

# LangChain imports with error handling
try:
    from langchain.agents import initialize_agent, Tool
    from langchain.agents.agent_types import AgentType
    from langchain_groq import ChatGroq
    # Prefer the official groq client if installed (new API). We'll detect and
    # create a small adapter to use the Responses API directly when present.
    try:
        import groq as _groq_client  # type: ignore
        GROQ_PY_AVAILABLE = True

        class GroqAdapter:
            """Minimal adapter to expose a callable that returns an object with
            `.content` similar to other LLM responses. This uses the Groq
            Responses API client when present.
            """
            def __init__(self, api_key: str, model: str):
                self._client = _groq_client.Client(api_key=api_key)
                self._model = model

            def __call__(self, prompt: str, **kwargs):
                """Call the installed groq client in a backward/forward-compatible way.

                Tries in order:
                - client.responses.create(...)
                - client.chat.completions.create(...)
                - any other reasonable completion endpoint if present
                Returns an object with a `.content` attribute (string).
                """
                resp = None
                # Try a few compatible client call shapes in order. We attempt
                # multiple candidate callables and multiple param signatures.
                tried = []
                def _try_call(fn, params):
                    try:
                        return fn(**params)
                    except TypeError:
                        # signature mismatch — bubble up to try other shapes
                        return None
                    except Exception:
                        return None

                # Inspect available attributes (helpful for debugging)
                try:
                    client_attrs = [a for a in dir(self._client) if not a.startswith('_')]
                except Exception:
                    client_attrs = []

                # Candidates: responses.create
                responses_obj = getattr(self._client, 'responses', None)
                if responses_obj is not None and hasattr(responses_obj, 'create'):
                    tried.append('responses.create')
                    resp = _try_call(responses_obj.create, {'model': self._model, 'input': prompt, **kwargs})

                # Candidate: client.chat may be an object or a callable factory
                if resp is None:
                    chat_obj = getattr(self._client, 'chat', None)
                    if callable(chat_obj):
                        try:
                            chat_obj = chat_obj()
                        except Exception:
                            pass
                    if chat_obj is not None:
                        # chat.completions.create
                        completions = getattr(getattr(chat_obj, 'completions', None), 'create', None)
                        if completions is not None:
                            tried.append('chat.completions.create')
                            # try a few param shapes
                            resp = _try_call(completions, {'model': self._model, 'input': prompt, **kwargs})
                            if resp is None:
                                resp = _try_call(completions, {'model': self._model, 'messages': [{'role':'user','content': prompt}], **kwargs})
                            if resp is None:
                                resp = _try_call(completions, {'model': self._model, 'prompt': prompt, **kwargs})

                # Candidate: top-level completions.create
                if resp is None:
                    comps = getattr(self._client, 'completions', None)
                    if comps is not None and hasattr(comps, 'create'):
                        tried.append('completions.create')
                        resp = _try_call(comps.create, {'model': self._model, 'input': prompt, **kwargs})
                        if resp is None:
                            resp = _try_call(comps.create, {'model': self._model, 'messages': [{'role':'user','content': prompt}], **kwargs})

                # If nothing worked, include attrs in the error for debugging
                if resp is None:
                    raise RuntimeError(f'Unable to call groq client: no compatible completion method found; tried={tried}; client_attrs={client_attrs}')

                # If still None, raise so callers see the underlying problem
                if resp is None:
                    raise RuntimeError('Unable to call groq client: no compatible completion method found')

                # Extract text from a variety of possible response shapes
                text = None
                try:
                    # Common convenience field used by some clients
                    if hasattr(resp, 'output_text') and resp.output_text:
                        text = resp.output_text
                    # Some clients expose `output` as a list/tree
                    if text is None:
                        out = getattr(resp, 'output', None)
                        if out:
                            # Drill into nested structures safely
                            # e.g. output[0].content[0].text OR output_text-like strings
                            first = out[0]
                            # message-like
                            msg = getattr(first, 'message', None)
                            if msg is not None:
                                # message.content may be string or list
                                cont = getattr(msg, 'content', None)
                                if isinstance(cont, str):
                                    text = cont
                                elif isinstance(cont, (list, tuple)) and cont:
                                    text = getattr(cont[0], 'text', str(cont[0]))
                            # direct text field
                            if text is None:
                                text = getattr(first, 'text', None) or getattr(first, 'content', None)
                                if isinstance(text, (list, tuple)):
                                    text = text[0] if text else ''
                    # Some clients use choices -> text
                    if text is None and hasattr(resp, 'choices'):
                        ch = getattr(resp, 'choices')
                        if isinstance(ch, (list, tuple)) and ch:
                            first = ch[0]
                            text = getattr(first, 'text', None) or getattr(first, 'message', None)
                            if hasattr(text, 'get') and isinstance(text, dict):
                                # e.g. {'content': '...'}
                                text = text.get('content') or str(text)
                    # Fallback to stringifying the object
                    if text is None:
                        text = str(resp)
                except Exception:
                    text = str(resp)

                # Ensure we return a plain string in .content so callers
                # won't accidentally display object reprs. If `text` is
                # already a message-like object with a `.content` field,
                # extract and coerce to string.
                try:
                    if not isinstance(text, str):
                        # If it's a message-like object, try to get .content
                        cont = getattr(text, 'content', None) or getattr(text, 'text', None)
                        if isinstance(cont, str):
                            text = cont
                        else:
                            # Last resort, stringify the original
                            text = str(text)
                except Exception:
                    text = str(text)

                return type('R', (), {'content': text})()

            # Compatibility method: many code paths expect an `invoke` method
            # that takes a prompt and returns an object with a `.content` attr.
            def invoke(self, prompt: str, **kwargs):
                return self.__call__(prompt, **kwargs)

            # Some code may expect batch/generate-like behavior. Provide a
            # minimal `generate` that accepts a list of prompts and returns a
            # list of response-like objects with `.content`.
            def generate(self, prompts, **kwargs):
                out = []
                try:
                    for p in prompts:
                        out.append(self.__call__(p, **kwargs))
                except Exception:
                    # Best-effort: if the client can handle batching natively,
                    # callers should replace this with a direct client call.
                    raise
                return out

    except Exception:
        _groq_client = None
        GROQ_PY_AVAILABLE = False
    from langchain_google_genai import ChatGoogleGenerativeAI
    from langchain_community.utilities import SerpAPIWrapper
    from langchain_community.document_loaders import PyPDFLoader, TextLoader
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    from langchain_huggingface import HuggingFaceEmbeddings
    from langchain_community.vectorstores import FAISS
    from langchain.chains import RetrievalQA
    from langchain_community.llms import HuggingFaceHub
    from langchain_core.documents import Document
    LANGCHAIN_AVAILABLE = True
except ImportError as e:
    # LangChain (or related) packages not available in this environment.
    # Define placeholders so the rest of the module can be imported and
    # runtime checks can decide whether to use these features.
    try:
        st.error(f"❌ LangChain components missing: {e}")
    except Exception:
        pass
    initialize_agent = None
    Tool = None
    AgentType = None
    ChatGroq = None
    ChatGoogleGenerativeAI = None
    SerpAPIWrapper = None
    PyPDFLoader = None
    TextLoader = None
    RecursiveCharacterTextSplitter = None
    HuggingFaceEmbeddings = None
    FAISS = None
    RetrievalQA = None
    HuggingFaceHub = None
    Document = None
    LANGCHAIN_AVAILABLE = False

# Voice imports with error handling
try:
    from gtts import gTTS
    import speech_recognition as sr
    from io import BytesIO
    import wave
    VOICE_AVAILABLE = True
    MICROPHONE_AVAILABLE = True
except ImportError as e:
    # Voice dependencies missing — create placeholders so module imports.
    VOICE_AVAILABLE = False
    MICROPHONE_AVAILABLE = False
    gTTS = None
    sr = None
    BytesIO = None
    wave = None
    # Try best-effort: TTS may still be available even if speech_recognition is not
    if "speech_recognition" not in str(e):
        try:
            from gtts import gTTS
            from io import BytesIO
            VOICE_AVAILABLE = True
        except Exception:
            gTTS = None
            BytesIO = None

# === Optional OCR / document parsing tools (best-effort imports) ===
try:
    import easyocr
    EASYOCR_AVAILABLE = True
except Exception:
    easyocr = None
    EASYOCR_AVAILABLE = False

try:
    import pytesseract
    PYTESSERACT_AVAILABLE = True
except Exception:
    pytesseract = None
    PYTESSERACT_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except Exception:
    fitz = None
    PYMUPDF_AVAILABLE = False

try:
    import docx
    PYDOCX_AVAILABLE = True
except Exception:
    docx = None
    PYDOCX_AVAILABLE = False

try:
    import pptx
    PYTTHON_PPTX_AVAILABLE = True
except Exception:
    pptx = None
    PYTTHON_PPTX_AVAILABLE = False

# Default OCR languages (comma-separated env var)
_OCR_LANGS = os.getenv('OCR_LANGS', 'en,hi').split(',')

def _ocr_image_bytes(raw_bytes, langs=None):
    """Extract text from image bytes using EasyOCR (preferred) then pytesseract fallback.
    Returns the extracted Unicode text.
    """
    langs = langs or _OCR_LANGS
    text_parts = []
    # Try EasyOCR first
    if EASYOCR_AVAILABLE:
        try:
            # initialize reader lazily for performance
            reader = easyocr.Reader(langs, gpu=False)
            import io
            from PIL import Image
            im = Image.open(io.BytesIO(raw_bytes)).convert('RGB')
            results = reader.readtext(im)
            # results: list of (bbox, text, confidence) or (bbox, text)
            for r in results:
                # r[1] is the text
                if isinstance(r, (list, tuple)) and len(r) >= 2:
                    text_parts.append(r[1])
            if text_parts:
                return '\n'.join(text_parts)
        except Exception:
            pass

    # Fallback to pytesseract if available
    if PYTESSERACT_AVAILABLE:
        try:
            from PIL import Image
            import io
            im = Image.open(io.BytesIO(raw_bytes)).convert('RGB')
            # pytesseract expects language codes like 'eng'/'hin' — default to eng
            t_lang = 'eng'
            if langs:
                # map first lang heuristically
                first = langs[0].strip()
                if first == 'hi':
                    t_lang = 'hin'
            return pytesseract.image_to_string(im, lang=t_lang)
        except Exception:
            pass

    # Last resort: return empty string
    return ''

# === Configuration ===
load_dotenv()

# Create directories
CACHE_DIR = Path("cache")
LOGS_DIR = Path("logs") 
UPLOADS_DIR = Path("uploads")
for dir_path in [CACHE_DIR, LOGS_DIR, UPLOADS_DIR]:
    dir_path.mkdir(exist_ok=True)

# --- Logging helpers ---
def _mask(key: str) -> str:
    if not key:
        return ""
    if len(key) <= 8:
        return "***"
    return f"{key[:4]}***{key[-4:]}"

def _latest_log_path() -> Path:
    LOGS_DIR.mkdir(exist_ok=True)
    # ChatbotLogger in logger.py writes daily files like logs/chatbot_YYYYMMDD.log
    # We don't know the exact name from here, but we can pick the newest .log
    logs = sorted(LOGS_DIR.glob("*.log"), key=lambda p: p.stat().st_mtime, reverse=True)
    return logs[0] if logs else LOGS_DIR / "chatbot.log"

# Initial startup log (ensures a log file exists)
try:
    logger.info("App startup: initializing components")
except Exception:
    pass

# === API Configuration ===
class APIConfig:
    """Centralized API configuration management"""
    
    def __init__(self):
        self.groq_key = os.getenv("GROQ_API_KEY")
        self.google_key = os.getenv("GOOGLE_API_KEY") 
        self.huggingface_token = os.getenv("HUGGINGFACE_API_TOKEN")
        self.serpapi_key = os.getenv("SERPAPI_API_KEY")
        try:
            logger.info(
                "API keys loaded",
                keys_present={
                    "groq": bool(self.groq_key),
                    "google": bool(self.google_key),
                    "huggingface": bool(self.huggingface_token),
                    "serpapi": bool(self.serpapi_key),
                },
            )
        except Exception:
            pass
    
    @property
    def has_groq(self):
        return bool(self.groq_key and len(self.groq_key) > 10)
    
    @property 
    def has_google(self):
        return bool(self.google_key and len(self.google_key) > 10)
    
    @property
    def has_huggingface(self):
        return bool(self.huggingface_token and len(self.huggingface_token) > 10)
    
    @property
    def has_serpapi(self):
        return bool(self.serpapi_key and len(self.serpapi_key) > 10)
    
    @property
    def available_providers(self):
        providers = []
        if self.has_groq:
            providers.append("groq")
        if self.has_google:
            providers.append("gemini")
        if self.has_huggingface:
            providers.append("huggingface")
        return providers
    
    def get_status_display(self):
        """Get formatted status for sidebar"""
        status = []
        status.append(f"{'✅' if self.has_groq else '❌'} Groq: {'Connected' if self.has_groq else 'Not configured'}")
        status.append(f"{'✅' if self.has_google else '❌'} Gemini: {'Connected' if self.has_google else 'Not configured'}")
        status.append(f"{'✅' if self.has_huggingface else '❌'} HuggingFace: {'Connected' if self.has_huggingface else 'Not configured'}")
        status.append(f"{'✅' if self.has_serpapi else '❌'} Web Search: {'Available' if self.has_serpapi else 'Unavailable'}")
        return status

config = APIConfig()

# === Avatar Management ===
class AvatarManager:
    """Manage user and bot avatars with fallbacks"""
    
    def __init__(self):
        self.assets_dir = Path("assets")
        self.user_avatar_path = self.assets_dir / "user_avatar.png"
        self.bot_avatar_path = self.assets_dir / "bot_avatar.png"
        self.user_avatar = None
        self.bot_avatar = None
        self.load_avatars()
    
    def load_avatars(self):
        """Load avatars from assets folder"""
        self.user_avatar = self.get_user_avatar()
        self.bot_avatar = self.get_bot_avatar()
    
    def get_user_avatar(self):
        """Get user avatar with fallback to emoji"""
        if self.user_avatar_path.exists():
            try:
                # Check if it's a valid image file
                with open(self.user_avatar_path, 'rb') as f:
                    # Try to read a few bytes to verify it's a file
                    f.read(10)
                return str(self.user_avatar_path.absolute())
            except Exception:
                pass
        return "👤"
    
    def get_bot_avatar(self):
        """Get bot avatar with fallback to emoji"""
        if self.bot_avatar_path.exists():
            try:
                # Check if it's a valid image file
                with open(self.bot_avatar_path, 'rb') as f:
                    # Try to read a few bytes to verify it's a file
                    f.read(10)
                return str(self.bot_avatar_path.absolute())
            except Exception:
                pass
        return "🤖"
    
    def get_avatars(self):
        """Get both avatars as tuple"""
        return self.user_avatar, self.bot_avatar

avatar_manager = AvatarManager()

# === Mode Management ===
class ModeManager:
    """Manage online/offline mode with smart fallbacks"""
    
    def __init__(self):
        if "mode" not in st.session_state:
            st.session_state.mode = "auto"  # auto, online, offline
        if "last_online_test" not in st.session_state:
            st.session_state.last_online_test = None
        if "connection_status" not in st.session_state:
            st.session_state.connection_status = "unknown"
    
    def test_connection(self, provider="groq"):
        """Test if online connection works with retries, timeout, and error classification.

        Returns True if any primary endpoint succeeds; on failure stores detailed diagnostics
        in st.session_state.connection_diagnostics for sidebar display.
        """
        import requests
        import os
        import socket
        import time

        # Ensure diagnostics container
        if "connection_diagnostics" not in st.session_state:
            st.session_state.connection_diagnostics = {}

        diagnostics = {
            "attempts": [],
            "provider": provider,
            "final": "fail",
            "started": datetime.utcnow().isoformat() + "Z",
            "proxies": {},
        }

        # Capture proxy environment
        for pv in ["HTTP_PROXY", "HTTPS_PROXY", "http_proxy", "https_proxy", "NO_PROXY", "no_proxy"]:
            val = os.getenv(pv)
            if val:
                diagnostics["proxies"][pv] = val

        if not config.available_providers:
            diagnostics["note"] = "No providers configured"
            st.session_state.connection_diagnostics = diagnostics
            try:
                logger.debug("Connection test: no providers configured")
            except Exception:
                pass
            return False

        # Base endpoints
        endpoints = [
            ("google_204", "https://www.google.com/generate_204"),
            ("httpbin", "https://httpbin.org/status/200"),
        ]
        if provider == "groq":
            endpoints.append(("groq", "https://api.groq.com"))
        elif provider == "gemini":
            endpoints.append(("googleapis", "https://generativelanguage.googleapis.com"))

        # Settings (with defaults; allow override via env)
        timeout = float(os.getenv("CONNECTION_TEST_TIMEOUT", "5"))  # seconds per attempt
        max_retries = int(os.getenv("CONNECTION_TEST_RETRIES", "2"))  # additional retries after first
        backoff_base = float(os.getenv("CONNECTION_TEST_BACKOFF", "0.5"))  # exponential backoff base seconds

        # TLS verification strategy
        ca_bundle = os.getenv("REQUESTS_CA_BUNDLE") or os.getenv("SSL_CERT_FILE")
        allow_insecure = os.getenv("ALLOW_INSECURE_SSL", "false").lower() == "true"
        # Use typing.Union for Python <3.10 compatibility
        verify: Union[bool, str] = ca_bundle if ca_bundle else (False if allow_insecure else True)

        def classify_error(exc: Exception) -> str:
            s = str(exc).lower()
            if any(k in s for k in ["ssl", "certificate", "handshake"]):
                return "ssl"
            if any(k in s for k in ["timeout", "timed out"]):
                return "timeout"
            if any(k in s for k in ["proxy", "tunnel"]):
                return "proxy"
            if any(k in s for k in ["name or service", "dns", "getaddrinfo"]):
                return "dns"
            if any(k in s for k in ["refused", "reset", "unreachable"]):
                return "network"
            return "other"

        # Attempt to capture DNS earlier
        try:
            socket.gethostbyname("www.google.com")
            diagnostics["dns_google"] = "ok"
        except Exception as e:
            diagnostics["dns_google"] = f"fail:{e.__class__.__name__}"
        success = False
        predominant_class = None
        for name, url in endpoints:
            attempt_index = 0
            while attempt_index <= max_retries:
                started_at = time.time()
                attempt_record = {
                    "endpoint": name,
                    "url": url,
                    "attempt": attempt_index + 1,
                    "verify": (verify if isinstance(verify, bool) else "custom_bundle"),
                }
                try:
                    resp = requests.get(url, timeout=timeout, verify=verify)
                    attempt_record["status_code"] = resp.status_code
                    attempt_record["elapsed_ms"] = int((time.time() - started_at) * 1000)
                    if resp.status_code < 500:
                        attempt_record["result"] = "ok"
                        diagnostics["attempts"].append(attempt_record)
                        success = True
                        break
                    else:
                        attempt_record["result"] = "http_error"
                except Exception as e:
                    attempt_record["error"] = str(e)[:400]
                    attempt_record["error_class"] = e.__class__.__name__
                    attempt_record["classification"] = classify_error(e)
                    attempt_record["elapsed_ms"] = int((time.time() - started_at) * 1000)
                diagnostics["attempts"].append(attempt_record)
                if attempt_record.get("result") == "ok":
                    break
                if not predominant_class and attempt_record.get("classification"):
                    predominant_class = attempt_record.get("classification")
                attempt_index += 1
                if attempt_index <= max_retries:
                    delay = backoff_base * (2 ** (attempt_index - 1))
                    time.sleep(min(delay, 3.0))
            if success:
                break

        diagnostics["final"] = "success" if success else "fail"
        diagnostics["predominant_error"] = predominant_class

        # Auto secondary test for SSL-only failures with custom bundle
        if (not success and ca_bundle and all(a.get("classification") == "ssl" for a in diagnostics["attempts"] if a.get("classification"))):
            secondary = {"variant": "system_trust_retry", "attempts": []}
            try:
                orig_ca = ca_bundle
                os.environ.pop("REQUESTS_CA_BUNDLE", None)
                os.environ.pop("SSL_CERT_FILE", None)
                for name, url in endpoints:
                    try:
                        resp = requests.get(url, timeout=timeout, verify=True)
                        secondary["attempts"].append({
                            "endpoint": name,
                            "status_code": resp.status_code,
                            "result": "ok" if resp.status_code < 500 else "http_error"
                        })
                        if resp.status_code < 500:
                            secondary["success"] = True
                    except Exception as e:
                        secondary["attempts"].append({
                            "endpoint": name,
                            "error": str(e)[:300],
                            "classification": classify_error(e)
                        })
                os.environ["REQUESTS_CA_BUNDLE"] = orig_ca  # restore
            except Exception as e:
                secondary["error"] = str(e)
            diagnostics["secondary"] = secondary

        # Update state and graceful degrade logic
        st.session_state.connection_diagnostics = diagnostics
        if success:
            st.session_state.connection_fail_streak = 0
        else:
            st.session_state.connection_fail_streak = st.session_state.get("connection_fail_streak", 0) + 1
            st.session_state.last_offline_reason = predominant_class or "unknown"
            threshold = int(os.getenv("AUTO_OFFLINE_FAIL_THRESHOLD", "3"))
            if st.session_state.mode == "auto" and st.session_state.connection_fail_streak >= threshold:
                # Force offline degrade
                st.session_state.connection_status = "offline"
                try:
                    logger.warning("Auto degraded to offline", fail_streak=st.session_state.connection_fail_streak, reason=st.session_state.last_offline_reason)
                except Exception:
                    pass

        try:
            logger.info("Connection test complete", success=success, attempts=len(diagnostics["attempts"]))
        except Exception:
            pass
        return success
    
    def get_current_mode(self):
        """Get effective current mode"""
        if st.session_state.mode == "offline":
            return "offline"
        elif st.session_state.mode == "online":
            return "online"
        else:  # auto mode
            # Test connection if not tested recently
            now = datetime.now()
            if (st.session_state.last_online_test is None or 
                (now - st.session_state.last_online_test).seconds > 30):
                
                st.session_state.connection_status = "testing"
                is_online = self.test_connection()
                st.session_state.connection_status = "online" if is_online else "offline"
                st.session_state.last_online_test = now
            
            return st.session_state.connection_status
    
    def force_mode(self, mode):
        """Force specific mode"""
        st.session_state.mode = mode
        st.session_state.last_online_test = None
    
    def set_mode(self, mode):
        """Set the current mode (alias for force_mode)"""
        self.force_mode(mode)

mode_manager = ModeManager()

# === Caching System ===
@st.cache_data(ttl=3600)
def get_cached_response(query_hash):
    """Get cached response if available"""
    cache_file = CACHE_DIR / f"response_{query_hash}.json"
    if cache_file.exists():
        try:
            with open(cache_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception:
            return None
    return None

def save_cached_response(query_hash, response):
    """Save response to cache"""
    cache_file = CACHE_DIR / f"response_{query_hash}.json"
    try:
        with open(cache_file, 'w', encoding='utf-8') as f:
            json.dump({
                "response": response,
                "timestamp": datetime.now().isoformat(),
                "cached": True
            }, f, ensure_ascii=False, indent=2)
    except Exception as e:
        st.warning(f"Could not cache response: {e}")

# === Enhanced LLM Management with Smart Fallback ===
@st.cache_resource
def create_llm_with_fallback(provider, model_name=None, force_offline=False):
    """Create LLM with intelligent fallback mechanism"""
    if force_offline or not LANGCHAIN_AVAILABLE:
        try:
            logger.info("LLM creation skipped", reason="force_offline_or_no_langchain", provider=provider)
        except Exception:
            pass
        return None
    
    # Get current mode
    current_mode = mode_manager.get_current_mode()
    
    if current_mode == "offline":
        try:
            logger.info("Offline mode active; not creating online LLM")
        except Exception:
            pass
        return None
    
    # Try primary provider
    try:
        logger.info("Attempting LLM creation", provider=provider, model=model_name or "default")
    except Exception:
        pass
    llm = try_create_llm(provider, model_name)
    if llm:
        try:
            logger.info("LLM created successfully", provider=provider)
        except Exception:
            pass
        return llm
    
    # Auto-fallback to other providers if in auto mode
    if st.session_state.mode == "auto":
        st.warning(f"⚠️ {provider.title()} failed, trying fallback providers...")
        try:
            logger.warning("Primary provider failed; trying fallbacks", provider=provider)
        except Exception:
            pass
        
        for fallback_provider in config.available_providers:
            if fallback_provider != provider:
                llm = try_create_llm(fallback_provider, None)
                if llm:
                    try:
                        logger.info("Switched to fallback provider", provider=fallback_provider)
                    except Exception:
                        pass
                    st.info(f"✅ Switched to {fallback_provider.title()}")
                    return llm
        
        # All providers failed, switch to offline mode
        st.error("❌ All AI providers failed, switching to offline mode")
        try:
            logger.error("All providers failed; switching offline")
        except Exception:
            pass
        mode_manager.force_mode("offline")
    
    return None

def try_create_llm(provider, model_name=None):
    """Try to create LLM for specific provider"""
    try:
        if provider == "groq" and config.has_groq:
            # Default model (legacy). We keep a migration map for Groq model IDs
            # to automatically switch deprecated IDs to recommended replacements.
            default_model = model_name or "llama-3.1-8b-instant"
            groq_migration_map = {
                # See Groq deprecations: https://console.groq.com/docs/deprecations
                # Llama3 -> Llama 3.x replacements
                "llama-3.1-8b-instant": "llama-3.1-8b-instant",
                "llama-3.3-70b-versatile": "llama-3.3-70b-versatile",
                # Whisper / speech models
                "whisper-large-v3-turbo": "whisper-large-v3-turbo",
                # Gemma / other older models -> recommended replacements
                "llama-3.1-8b-instant": "llama-3.1-8b-instant",
                # DeepSeek / Qwen mappings
                "llama-3.3-70b-versatile": "llama-3.3-70b-versatile",
                "deepseek-r1-distill-qwen-32b": "qwen/qwen3-32b",
                # Mistral / Qwen transitions
                "qwen/qwen3-32b": "qwen/qwen3-32b",
                "qwen/qwen3-32b": "qwen/qwen3-32b",
                # Moonshot / kimi
                "moonshotai/kimi-k2-instruct-0905-0905": "moonshotai/kimi-k2-instruct-0905-0905-0905",
                # Mixtral / misc preview -> production mappings
                "qwen/qwen3-32b": "qwen/qwen3-32b",
                # Llama guard
                "llama-guard-3-8b": "meta-llama/llama-guard-4-12b",
            }
            model = groq_migration_map.get(default_model, default_model)
            if model != default_model:
                try:
                    st.warning(f"Model '{default_model}' is deprecated; using recommended replacement '{model}'")
                except Exception:
                    pass
                try:
                    logger.info("Migrating deprecated Groq model", from_model=default_model, to_model=model)
                except Exception:
                    pass

            # If the official groq client is available, use an adapter to call
            # the Responses API directly. Otherwise, fall back to ChatGroq.
            if GROQ_PY_AVAILABLE:
                try:
                    return GroqAdapter(api_key=config.groq_key, model=model)
                except Exception as e:
                    try:
                        logger.warning("Groq adapter instantiation failed", error=str(e))
                    except Exception:
                        pass
                    # fall through to ChatGroq fallback

            if ChatGroq is not None:
                # Instantiate ChatGroq with compatibility for different client versions
                try:
                    return ChatGroq(api_key=config.groq_key, model=model, temperature=0.7, max_tokens=2048)
                except TypeError:
                    try:
                        return ChatGroq(api_key=config.groq_key, model_name=model, temperature=0.7, max_tokens=2048)
                    except TypeError:
                        # Fall back to positional args
                        return ChatGroq(config.groq_key, model, temperature=0.7, max_tokens=2048)
        
        elif provider == "gemini" and config.has_google and ChatGoogleGenerativeAI is not None:
            model = model_name or "gemini-pro"
            return ChatGoogleGenerativeAI(
                model=model,
                google_api_key=config.google_key,
                temperature=0.7
            )
        
        elif provider == "huggingface" and config.has_huggingface and HuggingFaceHub is not None:
            model = model_name or "microsoft/DialoGPT-medium"
            return HuggingFaceHub(
                repo_id=model,
                huggingfacehub_api_token=config.huggingface_token,
                model_kwargs={"temperature": 0.7, "max_length": 512}
            )
        
        return None
        
    except Exception as e:
        error_msg = str(e).lower()
        try:
            logger.error("LLM creation error", provider=provider, error=str(e), traceback=traceback.format_exc())
        except Exception:
            pass
        if any(term in error_msg for term in ['connection', 'network', 'ssl', 'timeout']):
            st.warning(f"🌐 Network issue with {provider.title()}")
        elif 'unauthorized' in error_msg or 'invalid' in error_msg:
            st.error(f"🔑 Invalid API key for {provider.title()}")
        else:
            st.warning(f"⚠️ {provider.title()} temporarily unavailable")
        return None

@st.cache_resource  
def create_llm(provider, model_name=None):
    """Legacy function for backwards compatibility"""
    return try_create_llm(provider, model_name)

# === Offline Fallback Bot ===
class OfflineBot:
    """Enhanced offline chatbot with better responses"""
    
    def __init__(self):
        self.responses = {
            'greetings': [
                "Hello! I'm running in offline mode with limited capabilities.",
                "Hi there! I'm in offline mode but I can still help with document search.",
                "Hey! I'm working offline right now but can assist with basic queries."
            ],
            'help': [
                "In offline mode, I can:\n• Search uploaded documents\n• Provide basic responses\n• Help with simple questions\n\nFor full AI capabilities, please configure your API keys.",
                "I'm in offline mode with these features:\n• Document text search\n• Basic conversation\n• File processing\n\nAdd API keys for advanced AI responses!"
            ],
            'default': [
                "I understand you're asking about that topic. I'm in offline mode with limited knowledge, but I can search any documents you've uploaded.",
                "That's an interesting question! I'm running offline, so I can't access external information, but I can help with document analysis if you upload files.",
                "I'd love to help with that! I'm in offline mode right now, but I can search through any documents you provide."
            ]
        }
    
    def invoke(self, prompt):
        """Generate contextual offline responses"""
        prompt_lower = prompt.lower()
        
        if any(word in prompt_lower for word in ['hello', 'hi', 'hey', 'greetings']):
            response = self.responses['greetings'][0]
        elif any(word in prompt_lower for word in ['help', 'what can you do', 'capabilities']):
            response = self.responses['help'][0]
        else:
            response = self.responses['default'][0]
        
        return type('Response', (), {'content': response})()

def _looks_like_pdf(header: bytes) -> bool:
    return header.startswith(b"%PDF")

def _looks_like_text(sample: bytes) -> bool:
    if not sample:
        return False
    try:
        text = sample.decode("utf-8", errors="ignore")
    except Exception:
        return False
    # Heuristic: > 85% printable characters
    printable = sum(1 for c in text if c.isprintable() or c in "\n\r\t")
    ratio = printable / max(1, len(text))
    return ratio >= 0.85

def _detect_file_kind(name: str, data: bytes) -> Optional[str]:
    n = (name or "").lower()
    head = data[:8]
    if n.endswith(".pdf") or _looks_like_pdf(head):
        return "pdf"
    if n.endswith(".txt") or _looks_like_text(data[:4096]):
        return "txt"
    # Common image extensions
    if n.endswith(('.png', '.jpg', '.jpeg', '.gif', '.tiff', '.bmp')):
        return "image"
    # OneNote types by extension (content is binary; we'll use best-effort extraction)
    if n.endswith(".one"):
        return "one"
    if n.endswith(".onepkg"):
        return "onepkg"
    if n.endswith(".onetoc2"):
        return "onetoc2"
    return None

def _strings_from_bytes(data: bytes, min_len: int = 5) -> str:
    try:
        text = data.decode("utf-8", errors="ignore")
    except Exception:
        return ""
    out = []
    buf = []
    for ch in text:
        if ch.isprintable() or ch in "\n\r\t":
            buf.append(ch)
        else:
            if len(buf) >= min_len:
                out.append("".join(buf))
            buf = []
    if len(buf) >= min_len:
        out.append("".join(buf))
    return "\n".join(out)


def get_response_text(llm_result: Any) -> str:
    """Safely extract text from various LLM response shapes.

    Accepts objects with `.content`, `.text`, `.get('content')`, or plain strings.
    """
    try:
        if llm_result is None:
            return ""
        # Some LLMs return an object with .content
        if hasattr(llm_result, 'content'):
            return str(llm_result.content)
        if hasattr(llm_result, 'text'):
            return str(llm_result.text)
        # Some return a dict-like
        if isinstance(llm_result, dict):
            for key in ('content', 'text', 'output'):
                if key in llm_result and llm_result[key]:
                    return str(llm_result[key])
        # Some return a plain string
        if isinstance(llm_result, str):
            return llm_result
        # Fallback: try stringifying
        return str(llm_result)
    except Exception:
        return ""

def _try_export_onenote_to_pdf(one_path: Path, pdf_path: Path) -> bool:
    """Best-effort export using OneNote COM (Windows, requires OneNote)."""
    try:
        import win32com.client  # type: ignore
        onenote = win32com.client.Dispatch("OneNote.Application")
        # OpenHierarchy returns (hierarchyId, outXML); Python COM often returns just id
        try:
            # Some installs: returns a tuple
            result = onenote.OpenHierarchy(str(one_path), "", 0)
            if isinstance(result, (list, tuple)):
                hierarchy_id = result[0]
            else:
                hierarchy_id = result
        except Exception:
            # Fallback: try without placeholders
            hierarchy_id = onenote.OpenHierarchy(str(one_path), None, 0)
        # PublishFormat: 3 is PDF per OneNote 2016 API
        onenote.Publish(hierarchy_id, str(pdf_path), 3)
        return pdf_path.exists() and pdf_path.stat().st_size > 0
    except Exception as e:
        try:
            logger.warning("OneNote COM export failed", error=str(e))
        except Exception:
            pass
        return False

# === Document Processing ===
def process_document(uploaded_file):
    """Wrapper for document processing that extracts raw bytes and filename
    and delegates to a cached helper which accepts primitive types (bytes,str).
    This avoids Streamlit cache hashing issues with file-like objects.
    """
    if not uploaded_file:
        return None

    try:
        raw = uploaded_file.getvalue() if hasattr(uploaded_file, "getvalue") else uploaded_file.getbuffer()
        raw_bytes = bytes(raw)
        filename = getattr(uploaded_file, 'name', 'uploaded') or 'uploaded'
        return _process_document_bytes(raw_bytes, filename)
    except Exception as e:
        try:
            st.error(f"❌ Error processing document: {e}")
        except Exception:
            pass
        try:
            logger.error("Document processing wrapper error", error=str(e))
        except Exception:
            pass
        return None


@st.cache_data
def _process_document_bytes(raw_bytes: bytes, filename: str):
    """Cached core document processing which operates on raw bytes and filename.
    Returns the same structure as the original `process_document`.
    """
    try:
        kind = _detect_file_kind(filename, raw_bytes)
        if not kind:
            try:
                logger.warning("Upload rejected: unknown kind", filename=filename, size=len(raw_bytes))
            except Exception:
                pass
            return None

        safe_name = filename
        if kind == "pdf" and not safe_name.lower().endswith(".pdf"):
            safe_name += ".pdf"
        if kind == "txt" and not safe_name.lower().endswith(".txt"):
            safe_name += ".txt"
        temp_path = UPLOADS_DIR / safe_name
        with open(temp_path, "wb") as f:
            f.write(raw_bytes)

        # Choose loader & splitter
        try:
            # Only use LangChain loaders for PDF or text files; images should go
            # through the manual/OCR fallback path to avoid trying to load binary
            # images with TextLoader.
            if (LANGCHAIN_AVAILABLE and PyPDFLoader and TextLoader and RecursiveCharacterTextSplitter
                    and kind in ("pdf", "txt")):
                if kind == "pdf":
                    loader = PyPDFLoader(str(temp_path))
                else:
                    loader = TextLoader(str(temp_path), encoding='utf-8')

                documents = loader.load()
                try:
                    temp_path.unlink(missing_ok=True)  # type: ignore[arg-type]
                except Exception:
                    pass

                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=1000,
                    chunk_overlap=200,
                    length_function=len
                )
                chunks_objs = text_splitter.split_documents(documents)
                chunk_texts = [c.page_content for c in chunks_objs]

                # If LangChain produced no chunks (e.g. scanned PDF with no extractable
                # text), fall back to manual extraction/OCR so we still return at least
                # one chunk.
                if not chunk_texts:
                    try:
                        raw_text = _strings_from_bytes(raw_bytes)
                    except Exception:
                        raw_text = ""

                    if not raw_text:
                        ext = temp_path.suffix.lower()
                        if ext == '.pdf' and PYMUPDF_AVAILABLE:
                            try:
                                texts = []
                                doc = fitz.open(str(temp_path))
                                for page in doc:
                                    try:
                                        pix = page.get_pixmap(dpi=150)
                                        img_bytes = pix.tobytes(output='png')
                                        t = _ocr_image_bytes(img_bytes)
                                        if t:
                                            texts.append(t)
                                    except Exception:
                                        continue
                                raw_text = "\n\n".join(texts)
                            except Exception:
                                raw_text = ""

                    # Ensure non-empty fallback
                    if not raw_text or len(raw_text.strip()) == 0:
                        raw_text = _strings_from_bytes(raw_bytes) or '(no extractable text found)'

                    # Create at least one chunk (plain strings)
                    chunk_size = 1000
                    overlap = 200
                    chunk_texts = []
                    i = 0
                    while i < len(raw_text):
                        part = raw_text[i:i+chunk_size]
                        chunk_texts.append(part)
                        i += (chunk_size - overlap)

                # Vector store (best effort)
                # Return simple, pickleable structures (lists of strings)
                return {
                    "type": "simple",
                    "chunks": chunk_texts,
                    "text_content": chunk_texts
                }
            else:
                # Simple fallback: try multiple extractors (text, office, OCR)
                try:
                    raw_text = _strings_from_bytes(raw_bytes)
                except Exception:
                    raw_text = ""

                # If no plain text found, try format-specific extraction
                if not raw_text:
                    ext = temp_path.suffix.lower()
                    # Office documents (docx)
                    if ext in ('.docx',) and PYDOCX_AVAILABLE:
                        try:
                            doc = docx.Document(str(temp_path))
                            raw_text = "\n".join([p.text for p in doc.paragraphs if p.text])
                        except Exception:
                            raw_text = ""
                    # PowerPoint
                    elif ext in ('.pptx',) and PYTTHON_PPTX_AVAILABLE:
                        try:
                            from pptx import Presentation
                            prs = Presentation(str(temp_path))
                            paras = []
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    try:
                                        if hasattr(shape, 'text') and shape.text:
                                            paras.append(shape.text)
                                    except Exception:
                                        continue
                            raw_text = "\n".join(paras)
                        except Exception:
                            raw_text = ""
                    # PDF: try to extract using PyMuPDF rendering + OCR if needed
                    elif ext == '.pdf' and PYMUPDF_AVAILABLE:
                        try:
                            texts = []
                            doc = fitz.open(str(temp_path))
                            for page in doc:
                                try:
                                    pix = page.get_pixmap(dpi=150)
                                    img_bytes = pix.tobytes(output='png')
                                    t = _ocr_image_bytes(img_bytes)
                                    if t:
                                        texts.append(t)
                                except Exception:
                                    continue
                            raw_text = "\n\n".join(texts)
                        except Exception:
                            raw_text = ""
                    # Images - attempt OCR
                    elif ext in ('.png', '.jpg', '.jpeg', '.gif', '.tiff', '.bmp'):
                        try:
                            raw_text = _ocr_image_bytes(raw_bytes)
                        except Exception:
                            raw_text = ""
                    else:
                        try:
                            raw_text = temp_path.read_text(encoding='utf-8', errors='ignore')
                        except Exception:
                            raw_text = ""

                # Remove temp file
                try:
                    temp_path.unlink(missing_ok=True)  # type: ignore[arg-type]
                except Exception:
                    pass

                # Ensure at least one chunk exists (avoid empty-chunk bug)
                if not raw_text or len(raw_text.strip()) == 0:
                    # Try extracting printable strings as a last resort
                    fallback = _strings_from_bytes(raw_bytes)
                    if fallback and fallback.strip():
                        raw_text = fallback
                    else:
                        # For images or binary PDFs where OCR is unavailable,
                        # provide a helpful placeholder indicating why no text
                        if kind == 'image':
                            raw_text = '(image file - no OCR backend installed or OCR returned empty)'
                        else:
                            raw_text = '(no extractable text found)'

                chunk_size = 1000
                overlap = 200
                chunk_texts = []
                i = 0
                while i < len(raw_text):
                    part = raw_text[i:i+chunk_size]
                    chunk_texts.append(part)
                    i += (chunk_size - overlap)

                return {
                    "type": "simple",
                    "chunks": chunk_texts,
                    "text_content": chunk_texts
                }
        except Exception as e:
            try:
                st.warning(f"Document processing fallback error: {e}")
            except Exception:
                pass
            return None
    except Exception as e:
        try:
            st.error(f"❌ Error processing document: {e}")
        except Exception:
            pass
        try:
            logger.error("Document processing error", error=str(e))
        except Exception:
            pass
        return None

# === Search Functions ===
def search_documents(query, doc_data):
    """Search documents with fallback methods"""
    if not doc_data:
        return "No documents available to search."
    
    query_lower = query.lower()
    
    if doc_data["type"] == "vectorstore":
        try:
            # Use vector search
            retriever = doc_data["store"].as_retriever(search_kwargs={"k": 3})
            relevant_docs = retriever.get_relevant_documents(query)
            if relevant_docs:
                return "\n\n".join([doc.page_content[:500] + "..." for doc in relevant_docs])
        except Exception:
            pass
    
    # Fallback to simple text search
    results = []
    for text in doc_data.get("text_content", []):
        if any(term in text.lower() for term in query_lower.split()):
            results.append(text[:400] + "...")
            if len(results) >= 3:
                break
    
    return "\n\n".join(results) if results else "No relevant information found in documents."

@st.cache_data
def web_search(query):
    """Web search with error handling"""
    if not config.has_serpapi:
        return "Web search unavailable - SerpAPI key not configured"
    
    try:
        if SerpAPIWrapper is None:
            return "Web search unavailable - SerpAPI wrapper not installed"
        search = SerpAPIWrapper(serpapi_api_key=config.serpapi_key)
        return search.run(query)
    except Exception as e:
        return f"Web search error: {e}"

# === Enhanced Voice Functions ===
def text_to_speech(text):
    """Convert text to speech with error handling"""
    if not VOICE_AVAILABLE:
        st.warning("Text-to-speech unavailable - install gtts package")
        return None
        
    try:
        if gTTS is None or BytesIO is None:
            st.warning("Text-to-speech backend not available")
            return None
        tts = gTTS(text=text[:500], lang='en')  # Limit length
        audio_buffer = BytesIO()
        tts.write_to_fp(audio_buffer)
        audio_buffer.seek(0)
        return audio_buffer
    except Exception as e:
        st.error(f"Text-to-speech error: {e}")
        return None

def speech_to_text():
    """Enhanced speech-to-text with better error handling"""
    if not MICROPHONE_AVAILABLE:
        return None, "🎤 Microphone not available - install pyaudio"
    if sr is None:
        return None, "Microphone backend not installed"

    try:
        r = sr.Recognizer()

        # Use default microphone
        with sr.Microphone() as source:
            st.info("🎤 Listening... Speak now!")

            # Adjust for ambient noise
            r.adjust_for_ambient_noise(source, duration=1)

            # Listen for audio
            audio = r.listen(source, timeout=10, phrase_time_limit=10)

        st.info("🔄 Processing speech...")

        # Try multiple recognition engines
        for engine, method in [
            ("Google", r.recognize_google),
            ("Sphinx", r.recognize_sphinx),
        ]:
            try:
                text = method(audio)
                if text.strip():
                    return text.strip(), f"✅ Recognized via {engine}"
            except AttributeError:
                # method may not exist for some engines
                continue
            except Exception:
                # Recognition failed for this engine, try next
                continue

        return None, "❌ Could not understand audio"

    except getattr(sr, 'WaitTimeoutError', Exception):
        return None, "⏰ No speech detected in 10 seconds"
    except Exception as e:
        return None, f"❌ Microphone error: {str(e)}"

def create_voice_input_component():
    """Create voice input UI component"""
    col1, col2 = st.columns([3, 1])
    
    with col2:
        if st.button("🎤 Voice Input", disabled=not MICROPHONE_AVAILABLE, help="Click to speak your message"):
            if MICROPHONE_AVAILABLE:
                with st.spinner("🎤 Listening..."):
                    text, status = speech_to_text()
                    if text:
                        st.session_state.voice_input = text
                        st.success(f"{status}: '{text}'")
                        st.rerun()
                    else:
                        st.error(status)
            else:
                st.error("Microphone not available - install pyaudio package")
    
    # Return voice input if available
    if hasattr(st.session_state, 'voice_input'):
        voice_text = st.session_state.voice_input
        del st.session_state.voice_input
        return voice_text
    
    return None

# === Main App Interface ===
def setup_page():
    """Configure Streamlit page"""
    st.set_page_config(
        page_title="🤖 AI Chatbot",
        page_icon="🤖",
        layout="wide",
        initial_sidebar_state="expanded"
    )

def show_welcome_screen():
    """Show welcome and setup screen"""
    st.title("🤖 Welcome to AI Chatbot")
    
    if not config.available_providers:
        st.warning("⚠️ No AI providers configured!")
        
        with st.expander("🔧 Quick Setup Guide", expanded=True):
            st.markdown("""
            ### Get started in 3 steps:
            
            1. **Create a `.env` file** in your project folder
            2. **Add your API keys:**
            ```
            GROQ_API_KEY=your_groq_key_here
            GOOGLE_API_KEY=your_google_key_here
            HUGGINGFACE_API_TOKEN=your_hf_token_here
            SERPAPI_API_KEY=your_serpapi_key_here
            ```
            3. **Restart the app**
            
            ### Get API Keys:
            - **Groq (Recommended)**: [console.groq.com](https://console.groq.com) - Free tier available
            - **Google Gemini**: [makersuite.google.com](https://makersuite.google.com) - Free tier available  
            - **HuggingFace**: [huggingface.co/settings/tokens](https://huggingface.co/settings/tokens) - Free
            - **SerpAPI**: [serpapi.com](https://serpapi.com) - Free tier available
            """)
        
        st.info("💡 The chatbot works in offline mode without API keys, but with limited capabilities.")
    
    else:
        st.success(f"✅ Ready to chat! {len(config.available_providers)} AI provider(s) available.")
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.metric("AI Providers", len(config.available_providers))
        
        with col2:
            st.metric("Web Search", "Available" if config.has_serpapi else "Unavailable")
            
        with col3:
            st.metric("Voice Features", "Available" if VOICE_AVAILABLE else "Limited")

def show_sidebar():
    """Enhanced sidebar with better organization"""
    with st.sidebar:
        st.title("🤖 AI Chatbot")

        # Provider selection
        if config.available_providers:
            provider = st.selectbox(
                "🧠 AI Provider:",
                config.available_providers,
                help="Choose your AI provider"
            )

            # Model selection
            if provider == "groq":
                model = st.selectbox("Model:", [
                    "llama-3.1-8b-instant",
                    "qwen/qwen3-32b",
                    "gemma-7b-it"
                ])
            elif provider == "gemini":
                model = st.selectbox("Model:", [
                    "gemini-pro",
                    "gemini-pro-vision"
                ])
            else:
                model = st.selectbox("Model:", [
                    "microsoft/DialoGPT-medium",
                    "facebook/blenderbot-400M-distill"
                ])
        else:
            provider = "offline"
            model = "offline"
            st.warning("⚠️ Running in offline mode")

        st.divider()

        # Document upload
        st.subheader("📄 Document Upload")
        uploaded_file = st.file_uploader(
            "Upload document:",
            help="Upload PDF or text files. OneNote users: export as PDF or paste text. Mislabeled files are auto-detected."
        )

        st.divider()

        # Mode Management
        st.subheader("🌐 Connection Mode")
        current_mode = mode_manager.get_current_mode()
        mode = st.selectbox(
            "Mode:",
            ["auto", "online", "offline"],
            index=["auto", "online", "offline"].index(current_mode),
            help="Auto: Smart switching, Online: Force online, Offline: Force offline"
        )

        if mode != current_mode:
            mode_manager.set_mode(mode)
            st.rerun()

        # Connection status
        if current_mode != "offline":
            connection_status = mode_manager.test_connection()
            status_color = "🟢" if connection_status else "🔴"
            st.write(f"{status_color} Connection: {'Online' if connection_status else 'Offline'}")

        st.divider()

        # Avatar Settings
        st.subheader("🎭 Avatar Settings")
        if st.button("🔄 Refresh Avatars"):
            avatar_manager.load_avatars()
            st.success("Avatars refreshed!")

        # Show current avatar status
        user_avatar, bot_avatar = avatar_manager.get_avatars()
        col1, col2 = st.columns(2)
        with col1:
            st.write("👤 User")
            user_is_custom = (
                isinstance(user_avatar, str)
                and (user_avatar.startswith("data:") or os.path.exists(user_avatar))
            )
            st.write("🖼️ Custom" if user_is_custom else "Default")
        with col2:
            st.write("🤖 Bot")
            bot_is_custom = (
                isinstance(bot_avatar, str)
                and (bot_avatar.startswith("data:") or os.path.exists(bot_avatar))
            )
            st.write("🖼️ Custom" if bot_is_custom else "Default")

        st.divider()

        # Settings
        st.subheader("⚙️ Settings")
        use_voice = st.checkbox(
            "🔊 Voice Features",
            value=False,
            disabled=not VOICE_AVAILABLE,
            help="Enable voice input and text-to-speech"
        )

        max_tokens = st.slider(
            "Response Length:",
            100, 2000, 500,
            help="Maximum response length"
        )

        # OCR language selector: allow user to change languages at runtime
        # Stored as both raw text and parsed list in session_state
        ocr_default_text = st.session_state.get('ocr_langs_text', ','.join(_OCR_LANGS))
        ocr_input = st.text_input(
            "OCR languages (comma-separated)",
            value=ocr_default_text,
            help="Comma-separated language codes (e.g. 'en,hi') used by OCR backends"
        )
        # Normalize and store in session state
        try:
            if ocr_input is not None:
                st.session_state['ocr_langs_text'] = ocr_input
                parsed = [s.strip() for s in ocr_input.split(',') if s.strip()]
                if parsed:
                    st.session_state['ocr_langs'] = parsed
        except Exception:
            pass

        st.divider()

        # Status
        st.subheader("📊 Status")
        for status_line in config.get_status_display():
            st.write(status_line)

        # Clear cache button
        if st.button("🗑️ Clear Cache"):
            st.cache_data.clear()
            st.cache_resource.clear()
            st.success("Cache cleared!")

        # Clear chat history button
        if st.button("🧹 Clear Chat History", help="Remove all chat messages from this session"):
            st.session_state.messages = []
            st.success("Chat history cleared.")
            st.rerun()

        st.divider()

        # Diagnostics
        st.subheader("🩺 Diagnostics")
        # Log level control (best-effort; logger levels are set at init)
        log_level = st.selectbox("Log level", ["DEBUG", "INFO", "WARNING", "ERROR"], index=1)
        try:
            # Update the underlying logger if possible
            logger.logger.setLevel(log_level)
        except Exception:
            pass

        # Connectivity tuning
        st.caption("Connection test settings")
        col_ct1, col_ct2, col_ct3 = st.columns(3)
        with col_ct1:
            timeout_val = st.number_input("Timeout (s)", min_value=1.0, max_value=30.0, value=float(os.getenv("CONNECTION_TEST_TIMEOUT", "5")), step=0.5, help="Per-attempt timeout")
        with col_ct2:
            retries_val = st.number_input("Retries", min_value=0, max_value=5, value=int(os.getenv("CONNECTION_TEST_RETRIES", "2")), step=1, help="Additional retries after first")
        with col_ct3:
            backoff_val = st.number_input("Backoff (s)", min_value=0.1, max_value=5.0, value=float(os.getenv("CONNECTION_TEST_BACKOFF", "0.5")), step=0.1, help="Initial exponential backoff")

        # Apply settings to environment for current process (ephemeral)
        if st.button("Apply test settings"):
            os.environ["CONNECTION_TEST_TIMEOUT"] = str(timeout_val)
            os.environ["CONNECTION_TEST_RETRIES"] = str(retries_val)
            os.environ["CONNECTION_TEST_BACKOFF"] = str(backoff_val)
            st.success("Updated connection test settings.")

        # Show certificate override status and provide quick controls
        st.caption("TLS trust configuration")
        ca_bundle = os.getenv("REQUESTS_CA_BUNDLE")
        ssl_cert_file = os.getenv("SSL_CERT_FILE")
        st.text(f"REQUESTS_CA_BUNDLE = {ca_bundle or '(not set)'}")
        st.text(f"SSL_CERT_FILE = {ssl_cert_file or '(not set)'}")
        ca_col1, ca_col2 = st.columns(2)
        with ca_col1:
            if st.button("Use system trust (clear overrides)"):
                # backup once per session
                if "_ca_backup" not in st.session_state:
                    st.session_state._ca_backup = {
                        "REQUESTS_CA_BUNDLE": ca_bundle,
                        "SSL_CERT_FILE": ssl_cert_file,
                    }
                try:
                    os.environ.pop("REQUESTS_CA_BUNDLE", None)
                    os.environ.pop("SSL_CERT_FILE", None)
                    st.success("Cleared CA overrides for this process. Re-check connectivity below.")
                except Exception as e:
                    st.error(f"Could not clear overrides: {e}")
        with ca_col2:
            if st.button("Restore overrides"):
                try:
                    backup = st.session_state.get("_ca_backup")
                    if backup:
                        if backup.get("REQUESTS_CA_BUNDLE"):
                            os.environ["REQUESTS_CA_BUNDLE"] = backup["REQUESTS_CA_BUNDLE"]
                        else:
                            os.environ.pop("REQUESTS_CA_BUNDLE", None)
                        if backup.get("SSL_CERT_FILE"):
                            os.environ["SSL_CERT_FILE"] = backup["SSL_CERT_FILE"]
                        else:
                            os.environ.pop("SSL_CERT_FILE", None)
                        st.success("Restored previous CA override values.")
                    else:
                        st.info("No backup available in this session.")
                except Exception as e:
                    st.error(f"Could not restore overrides: {e}")

        if st.button("🔎 Re-check connectivity (verbose)"):
            with st.spinner("Testing connection..."):
                ok = mode_manager.test_connection()
                st.success("Online") if ok else st.error("Offline")

        # Inline validation & self-test tools
        st.caption("Validation & provider self-test")
        col_v1, col_v2, col_v3 = st.columns(3)
        with col_v1:
            run_validator = st.button("🧪 Validate", help="Run project validator (compiles, imports, connectivity)")
        with col_v2:
            run_selftest = st.button("🧬 Self-Test", help="Run provider self-test (instantiate each provider)")
        with col_v3:
            run_both = st.button("🧪+🧬 Both", help="Run validator + self-test")

        def _run_validator():
            import subprocess
            import json
            import sys
            try:
                result = subprocess.run([sys.executable, "-X", "utf8", "tools/validate_project.py"], capture_output=True, text=True, timeout=120)
                out = result.stdout.strip()
                # Attempt parse; fallback to raw text
                try:
                    parsed = json.loads(out)
                except Exception:
                    parsed = {"raw": out[-5000:]}
                if result.returncode != 0:
                    parsed["returncode"] = str(result.returncode)
                return parsed
            except Exception as e:
                return {"error": str(e)}

        def _run_selftest():
            import subprocess
            import json
            import sys
            try:
                result = subprocess.run([sys.executable, "-X", "utf8", "tools/self_test.py"], capture_output=True, text=True, timeout=90)
                out = result.stdout.strip()
                try:
                    parsed = json.loads(out)
                except Exception:
                    parsed = {"raw": out[-4000:]}
                if result.returncode != 0:
                    parsed["returncode"] = str(result.returncode)
                return parsed
            except Exception as e:
                return {"error": str(e)}

        if run_validator or run_both:
            with st.spinner("Running validator..."):
                val = _run_validator()
            if not val.get("error"):
                st.success("Validator complete")
            else:
                st.error("Validator error")
            with st.expander("Validator Output", expanded=False):
                st.json(val)

        if run_selftest or run_both:
            with st.spinner("Running provider self-test..."):
                ste = _run_selftest()
            if not ste.get("error"):
                st.success("Self-test complete")
            else:
                st.error("Self-test error")
            with st.expander("Self-Test Output", expanded=False):
                st.json(ste)

        # Show last diagnostics summary
        diag = st.session_state.get("connection_diagnostics")
        if diag:
            with st.expander("Last connection diagnostics", expanded=False):
                st.json({
                    k: diag[k] for k in [
                        "provider", "final", "started"
                    ] if k in diag
                })
                st.write(f"Attempts: {len(diag.get('attempts', []))}")
                # Show condensed error classifications
                classes = {}
                for a in diag.get("attempts", []):
                    c = a.get("classification")
                    if c:
                        classes[c] = classes.get(c, 0) + 1
                if classes:
                    st.write("Error classes:", classes)
                if diag.get("secondary"):
                    sec = diag["secondary"]
                    st.write("Secondary (system trust) retry:")
                    st.json(sec)

        # Show last lines of the latest log
        log_file = _latest_log_path()
        st.caption(f"Log file: {log_file.name}")
        try:
            if log_file.exists():
                lines = log_file.read_text(encoding="utf-8", errors="ignore").splitlines()[-200:]
                st.code("\n".join(lines) if lines else "(log empty)", language="text")
            else:
                st.info("No log file yet. Interact with the app to generate logs.")
        except Exception as e:
            st.warning(f"Could not read log: {e}")

    return provider, model, uploaded_file, use_voice, max_tokens

def main():
    """Main application"""
    setup_page()
    
    # Initialize session state
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "doc_data" not in st.session_state:
        st.session_state.doc_data = None
    if "show_welcome" not in st.session_state:
        st.session_state.show_welcome = True
    
    # Show sidebar
    provider, model, uploaded_file, use_voice, max_tokens = show_sidebar()
    
    # Main content area
    if st.session_state.show_welcome and not st.session_state.messages:
        show_welcome_screen()
        if st.button("🚀 Start Chatting"):
            st.session_state.show_welcome = False
            st.rerun()
        return
    
    st.title("💬 Chat")
    
    # Process uploaded document
    if uploaded_file and uploaded_file != st.session_state.get("last_file"):
        with st.spinner("📖 Processing document..."):
            doc_data = process_document(uploaded_file)
            if doc_data:
                st.session_state.doc_data = doc_data
                st.session_state.last_file = uploaded_file
                st.success(f"✅ Document processed: {uploaded_file.name}")
            else:
                st.error("❌ Failed to process document")
    
    # Display chat history with custom avatars
    user_avatar, bot_avatar = avatar_manager.get_avatars()
    for message in st.session_state.messages:
        avatar = user_avatar if message["role"] == "user" else bot_avatar
        with st.chat_message(message["role"], avatar=avatar):
            st.markdown(message["content"])
            if message.get("cached"):
                st.caption("📚 From cache")
    
    # Voice input component
    voice_input = None 
    if use_voice:
        voice_input = create_voice_input_component()
    
    # Chat input - both text and voice
    prompt = None
    
    # Get text input
    text_input = st.chat_input("What would you like to know?")
    
    # Process voice input if available
    if voice_input:
        prompt = voice_input
        st.success(f"🎤 Voice input: {prompt}")
    elif text_input:
        prompt = text_input
    
    if prompt:
        # Get custom avatars
        user_avatar, bot_avatar = avatar_manager.get_avatars()
        
        # Add user message
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user", avatar=user_avatar):
            st.markdown(prompt)
        
        # Generate response
        with st.chat_message("assistant", avatar=bot_avatar):
            with st.spinner("🤔 Thinking..."):
                try:
                    try:
                        logger.info(
                            "Handling prompt",
                            provider=provider,
                            model=model,
                            mode=mode_manager.get_current_mode(),
                            has_doc=bool(st.session_state.doc_data),
                        )
                    except Exception:
                        pass
                    # Check cache first
                    query_hash = hashlib.md5(prompt.encode()).hexdigest()
                    cached = get_cached_response(query_hash)
                    
                    if cached:
                        response = cached["response"]
                        st.markdown(response)
                        st.caption("📚 From cache")
                        
                    else:
                        # Create LLM with smart fallback
                        force_offline = (provider == "offline" or mode_manager.get_current_mode() == "offline")
                        llm = create_llm_with_fallback(provider, model, force_offline)
                        
                        if not llm:
                            llm = OfflineBot()
                            try:
                                logger.warning("Using OfflineBot fallback", reason="llm_unavailable", provider=provider)
                            except Exception:
                                pass
                            if not force_offline:
                                st.info("🔄 Using offline mode")
                        
                        # Generate response
                        if st.session_state.doc_data:
                            # Include document search
                            doc_results = search_documents(prompt, st.session_state.doc_data)
                            
                            if isinstance(llm, OfflineBot):
                                response = f"{get_response_text(llm.invoke(prompt))}\n\n📄 **From your document:**\n{doc_results}"
                            else:
                                enhanced_prompt = f"{prompt}\n\nRelevant document content:\n{doc_results}"
                                response = get_response_text(llm.invoke(enhanced_prompt))
                        else:
                            # Regular response
                            response = get_response_text(llm.invoke(prompt))
                        
                        # Display response
                        st.markdown(response)
                        
                        # Save to cache
                        save_cached_response(query_hash, response)
                        
                        # Text-to-speech
                        if use_voice and response:
                            audio = text_to_speech(response)
                            if audio:
                                st.audio(audio, format='audio/mp3')
                    
                    # Add to session state
                    st.session_state.messages.append({
                        "role": "assistant", 
                        "content": response,
                        "cached": bool(cached)
                    })
                    
                except Exception as e:
                    error_response = f"I apologize, but I encountered an error: {str(e)}"
                    try:
                        logger.error("Chat handling error", error=str(e), traceback=traceback.format_exc())
                    except Exception:
                        pass
                    st.error(error_response)
                    st.session_state.messages.append({
                        "role": "assistant", 
                        "content": error_response
                    })

if __name__ == "__main__":
    main()
