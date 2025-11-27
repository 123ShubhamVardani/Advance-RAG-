"""Generate a project architecture PowerPoint deck.

Usage:
    python tools/generate_presentation.py

If python-pptx is not installed, install with:
    pip install python-pptx

Outputs:
    dist/AI_Chatbot_Architecture.pptx
"""
from pathlib import Path
import sys
from datetime import datetime

try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore
except ImportError as e:
    print("[generate_presentation] Import error:", e)
    print("sys.executable:", sys.executable)
    print("sys.path:\n" + "\n".join(sys.path))
    print("python-pptx not installed in this environment. Install with: pip install python-pptx")
    sys.exit(1)

DECK_NAME = "AI_Chatbot_Architecture"
OUTPUT_DIR = Path("dist")
OUTPUT_DIR.mkdir(exist_ok=True)

prs = Presentation()

# Helper styles
TITLE_SIZE = Pt(40)
SUBTITLE_SIZE = Pt(20)
BODY_SIZE = Pt(18)


def add_title_slide(title: str, subtitle: str = ""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_bullet_slide(title: str, bullets: list[str]):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            body.text = b
        else:
            p = body.add_paragraph()
            p.text = b
            p.level = 0


def add_content_slide(title: str, content: str):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for line in content.strip().splitlines():
        if not line.strip():
            p = tf.add_paragraph(); p.text = ""; continue
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
    # Remove the first empty default paragraph
    if tf.paragraphs and not tf.paragraphs[0].text:
        del tf.paragraphs[0]


add_title_slide("AI Chatbot Architecture", "Generated: " + datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"))

add_bullet_slide("Purpose", [
    "Multi-provider AI assistant (Groq, Gemini, HuggingFace)",
    "Document Q&A with RAG (PDF/TXT)",
    "Resilient offline fallback & diagnostics",
    "Corporate network friendly (proxy/SSL classification)"
])

add_bullet_slide("Core Pillars", [
    "Streamlit UI + Sidebar Diagnostics",
    "Mode Manager (auto/online/offline)",
    "Document Processing -> Embeddings -> Retrieval",
    "Provider Fallback Chain + OfflineBot",
])

add_content_slide("High-Level Flow", """
User Input / Upload
 -> Connectivity Test & Mode Decision
 -> (Optional) Document Chunking + Embeddings (FAISS)
 -> Retrieval adds context to prompt
 -> Provider selection (primary -> fallback chain)
 -> LLM or OfflineBot response
 -> Cache result + optional TTS
 -> Diagnostics & logs update
""")

add_bullet_slide("Resilience & Diagnostics", [
    "Connectivity classification: ssl / timeout / proxy / dns / network / other",
    "Variants: env vs system vs insecure",
    "Auto offline degrade after fail streak",
    "Validator + provider self-test buttons"
])

add_bullet_slide("Security & Trust", [
    "Conditional insecure SSL toggle (env var)",
    "System trust fallback (certifi_win32)",
    "CA override detection & secondary retry",
    "No keys stored in code (.env only)"
])

add_bullet_slide("Performance & UX", [
    "Response caching (hash-based)",
    "Configurable timeouts / retries / backoff",
    "Lightweight offline mode responses",
    "Voice input + TTS (optional)"
])

add_bullet_slide("Key Files", [
    "app.py - Main orchestrator",
    "tools/validate_project.py - Health & root cause",
    "tools/self_test.py - Provider smoke test",
    "CONNECTIVITY.md / TROUBLESHOOTING.md - Docs",
    "Dockerfile - Container packaging"
])

add_content_slide("Extension Ideas", """
Add new provider adapters
Integrate vector store persistence
Add auth / multi-user sessions
Export transcripts & analytics
CI pipeline to run validator
""")

outfile = OUTPUT_DIR / f"{DECK_NAME}.pptx"
prs.save(outfile)
print(f"✅ Presentation written to {outfile}")
